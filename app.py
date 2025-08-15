from flask import Flask, request, jsonify, render_template, send_from_directory
from flask import Flask, request, jsonify, send_from_directory, render_template
from flask_cors import CORS
import os
import json
import uuid
import logging
import time
import pdfplumber
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import requests
from sentence_transformers import SentenceTransformer
import io
import PyPDF2
import re
import json
import pdfplumber
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import requests
from sklearn.metrics.pairwise import cosine_similarity

# Initialize Flask app
app = Flask(__name__)
CORS(app)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['IMAGE_FOLDER'] = os.path.join(app.config['UPLOAD_FOLDER'], 'images')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB limit
ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx', 'txt'}

# Configuration
UPLOAD_FOLDER = 'uploads'
IMAGE_FOLDER = 'uploads/images'
DATA_FILE = 'document_data.json'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(IMAGE_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['IMAGE_FOLDER'] = IMAGE_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload size

# Constants
OLLAMA_API = "http://localhost:11434/api"
ALLOWED_EXTENSIONS = {'pdf'}
SIMILARITY_THRESHOLD = 0.75
TOP_K = 3
HEADING_SIMILARITY_THRESHOLD = 0.65
# Initialize SentenceTransformer model
print("Initializing SentenceTransformer model...")
model = SentenceTransformer('all-MiniLM-L6-v2')
print("Model initialized successfully")

# Debugging flag
DEBUG = True
# Load or initialize document data
DOCUMENT_DATA_FILE = 'document_data.json'

# Initialize components
document_chunks = {}
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
def load_document_data():
    print("Loading document_data.json...")
    try:
        with open(DOCUMENT_DATA_FILE, 'r') as f:
            data = json.load(f)
        print("document_data.json loaded successfully")
        return data
    except FileNotFoundError:
        print("document_data.json not found, initializing empty data")
        return {"documents": []}

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    filename='app.log'
)
logger = logging.getLogger(__name__)
def save_document_data(data):
    print("Saving document_data.json...")
    with open(DOCUMENT_DATA_FILE, 'w') as f:
        json.dump(data, f, indent=4)
    print("document_data.json saved successfully")

def log_debug(message):
    """Helper function for consistent debug logging"""
    if DEBUG:
        logger.debug(message)
        print(f"[DEBUG] {message}")
document_data = load_document_data()

# Load existing data
if os.path.exists(DATA_FILE):
    with open(DATA_FILE, 'r') as f:
        document_chunks = json.load(f)
    # Convert stored embeddings back to NumPy arrays
    for doc in document_chunks.values():
        if "chunk_embeddings" in doc:
            doc["chunk_embeddings"] = [np.array(emb) for emb in doc["chunk_embeddings"]]
    log_debug(f"Loaded {len(document_chunks)} documents from storage")
# Ensure upload and image folders exist
print("Creating upload and image folders if they don't exist...")
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['IMAGE_FOLDER'], exist_ok=True)
print("Folders created successfully")

# Helper Functions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_images_from_pdf(filepath):
    """Extract images from PDF and associate with nearby headings using PyMuPDF"""
def extract_text_from_pdf(file_path):
    text = ""
    headings = []
    print("Extracting text and headings from PDF...")
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                words = page.extract_words()
                for word in words:
                    if word['text'].isupper() or (word['text'][0].isupper() and len(word['text']) > 3):
                        headings.append({"text": word['text'], "page": page_num})
        print(f"PDF headings retrieved: {len(headings)} headings")
    except Exception as e:
        print(f"Error extracting PDF text/headings: {str(e)}")
        print("No headings found for PDF")
    return text, headings

def extract_images_from_pdf(file_path):
    images = []
    print("Extracting images from PDF...")
    try:
        doc = fitz.open(filepath)
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            # Get text blocks to identify headings
            text_blocks = page.get_text("dict")["blocks"]
            headings = []
            for block in text_blocks:
                if block["type"] == 0:  # Text block
                    for line in block["lines"]:
                        text = "".join(span["text"] for span in line["spans"]).strip()
                        # Simple heading detection (can be enhanced)
                        if text and re.match(r"^[A-Z][A-Za-z0-9 \-:]+$", text):
                            headings.append({"text": text, "y0": line["bbox"][1]})
            
            # Extract images
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = f"{uuid.uuid4()}.{image_ext}"
                image_filename = f"pdf_image_{uuid.uuid4()}.{image_ext}"
                image_path = os.path.join(app.config['IMAGE_FOLDER'], image_filename)
                
                # Save image to disk
                with open(image_path, "wb") as f:
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                # Find nearest heading (based on y-coordinate)
                image_rect = page.get_image_bbox(img)
                image_y = image_rect.y0
                nearest_heading = None
                min_distance = float('inf')
                for heading in headings:
                    distance = abs(heading["y0"] - image_y)
                    if distance < min_distance:
                        min_distance = distance
                        nearest_heading = heading["text"]
                
                # Store image metadata
                images.append({
                    "filename": image_filename,
                    "heading": nearest_heading or "No heading",
                    "page": page_num + 1
                })
        
        doc.close()
        return images
                heading = f"Image from page {page_num + 1}"
                images.append({"filename": image_filename, "heading": heading, "page": page_num + 1})
        print(f"Extracted {len(images)} images from PDF")
    except Exception as e:
        logger.error(f"Image extraction failed: {e}")
        return []

def extract_headings_and_content(pdf_url=None, raw_text=None, filepath=None, heading_pattern=r"^(#+)?\s*([A-Z][A-Za-z0-9 \-:]+)$"):
    """
    Enhanced heading extraction with image association
    """
    chunks = {}
    current_heading = None
    current_content = []
        print(f"Error extracting PDF images: {str(e)}")
    return images

    # Get text from either PDF, raw text, or filepath
def extract_text_from_pptx(file_path):
    text = ""
    images = []
    if filepath:
        try:
            with pdfplumber.open(filepath) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            images = extract_images_from_pdf(filepath)
        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            return chunks, images
    elif pdf_url:
        try:
            response = requests.get(pdf_url, timeout=20)
            pdf_stream = io.BytesIO(response.content)
            with pdfplumber.open(pdf_stream) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            return chunks, images
    elif raw_text:
        text = raw_text
    else:
        return chunks, images

    log_debug(f"Raw text sample:\n{text[:500]}...")
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue

        heading_match = re.match(heading_pattern, line)
        if heading_match:
            if current_heading:
                chunks[current_heading] = "\n".join(current_content).strip()
            current_heading = heading_match.group(2) if heading_match.group(2) else line
            current_content = []
            log_debug(f"Found heading: {current_heading}")
        elif current_heading:
            current_content.append(line)
    headings = []
    print("Extracting text and headings from PPTX...")
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                    if shape.is_placeholder and shape.placeholder_format.type in (1, 2):  # Title or Subtitle
                        headings.append({"text": shape.text, "slide": slide_num})
            text += slide_text + "\n"
        print(f"PPTX headings retrieved: {len(headings)} headings")
    except Exception as e:
        print(f"Error extracting PPTX text/headings: {str(e)}")
        print("No headings found for PPTX")
    return text, headings

    if current_heading:
        chunks[current_heading] = "\n".join(current_content).strip()
def extract_images_from_pptx(file_path):
    images = []
    print("Extracting images from PPTX...")
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = f"pptx_image_{uuid.uuid4()}.{image_ext}"
                    image_path = os.path.join(app.config['IMAGE_FOLDER'], image_filename)
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    heading = f"Image from slide {slide_num}"
                    images.append({"filename": image_filename, "heading": heading, "slide": slide_num})
        print(f"Extracted {len(images)} images from PPTX")
    except Exception as e:
        print(f"Error extracting PPTX images: {str(e)}")
    return images

    log_debug(f"Extracted {len(chunks)} headings, {len(images)} images")
    return chunks, images
def extract_text_from_docx(file_path):
    text = ""
    headings = []
    print("Extracting text and headings from DOCX...")
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                headings.append({"text": para.text, "level": para.style.name})
            text += para.text + "\n"
        print(f"DOCX headings retrieved: {len(headings)} headings")
    except Exception as e:
        print(f"Error extracting DOCX text/headings: {str(e)}")
        print("No headings found for DOCX")
    return text, headings

def extract_text_from_pdf(filepath):
    """Extract text from PDF with error handling"""
def extract_text_from_txt(file_path):
    text = ""
    headings = []
    print("Extracting text from TXT...")
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            return " ".join(page.extract_text() for page in reader.pages if page.extract_text())
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        print("No headings found for TXT (plain text)")  # TXT has no heading structure
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return ""
        print(f"Error extracting TXT text: {str(e)}")
    return text, headings

def chunk_text(text, chunk_size=1000):
    """Create context-preserving chunks with headings"""
    heading_content = extract_headings_and_content(raw_text=text)[0]
def chunk_text(text, max_chunk_size=500):
    print("Chunking text...")
    words = text.split()
    chunks = []
    
    if not heading_content:
        log_debug("No headings found, using standard chunking")
        sections = re.split(r'\n\s*\n', text)
        current_chunk = ""
        for section in sections:
            if len(current_chunk) + len(section) < chunk_size:
                current_chunk += "\n\n" + section
            else:
                chunks.append(current_chunk.strip())
                current_chunk = section
        if current_chunk:
            chunks.append(current_chunk.strip())
        return chunks
    
    for heading, content in heading_content.items():
        chunk = f"## {heading}\n{content}"
        while len(chunk) > chunk_size:
            split_pos = max(
                chunk.rfind('\n\n', 0, chunk_size),
                chunk.rfind('. ', 0, chunk_size)
            )
            if split_pos == -1:
                split_pos = chunk_size
            chunks.append(chunk[:split_pos].strip())
            chunk = f"## {heading}\n" + chunk[split_pos:].lstrip()
        chunks.append(chunk.strip())
    current_chunk = []
    current_length = 0
    for word in words:
        current_length += len(word) + 1
        if current_length > max_chunk_size:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word) + 1
        else:
            current_chunk.append(word)
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    print(f"Text split into {len(chunks)} chunks")
    return chunks

def get_relevant_chunks(query, chunks, chunk_embeddings):
    """Get most relevant chunks with similarity scores using precomputed embeddings"""
    if not chunks or not chunk_embeddings:
        log_debug("⚠️ Warning: No chunks or embeddings available for similarity calculation")
        return []
    
    query_embed = embedding_model.encode([query])
    chunk_embeds = np.array(chunk_embeddings)
    
    if len(query_embed.shape) == 1:
        query_embed = query_embed.reshape(1, -1)
    if len(chunk_embeds.shape) == 1:
        chunk_embeds = chunk_embeds.reshape(1, -1)
    
    try:
        sims = cosine_similarity(query_embed, chunk_embeds)[0]
        top_indices = np.argsort(sims)[-TOP_K:][::-1]
        return [(chunks[i], sims[i]) for i in top_indices]
    except ValueError as e:
        log_debug(f"⚠️ Similarity calculation failed: {str(e)}")
        return []

def query_llama(prompt, model="llama3"):
    """Query LLM with debug logging"""
def generate_embeddings(chunks):
    print("Generating embeddings for chunks...")
    try:
        log_debug(f"Sending to LLM: {prompt[:200]}...")
        start_time = time.time()
        
        response = requests.post(
            f"{OLLAMA_API}/generate",
            json={
                "model": model,
                "prompt": prompt,
                "stream": False,
                "options": {"temperature": 0.3}
            },
            timeout=60
        )
        response.raise_for_status()
        
        result = response.json().get("response", "").strip()
        elapsed = time.time() - start_time
        log_debug(f"LLM response ({elapsed:.2f}s): {result[:200]}...")
        return result
        embeddings = model.encode(chunks, convert_to_tensor=False)
        print("Embeddings generated successfully")
        return embeddings
    except Exception as e:
        logger.error(f"LLM query failed: {e}")
        return "I encountered an error processing your request."
        print(f"Error generating embeddings: {str(e)}")
        return []

# Routes
@app.route('/upload', methods=['POST'])
def upload_files():
def upload_file():
    print("Received file upload request")
    if 'files[]' not in request.files:
        print("No files uploaded")
        return jsonify({"success": False, "message": "No files uploaded"}), 400
    
    files = request.files.getlist('files[]')
    uploaded_docs = []
    errors = []
    
    uploaded_documents = []
    for file in files:
        if not file or file.filename == '':
            errors.append(f"Empty SUCCESS file part for one of the files")
            continue
            
        if not allowed_file(file.filename):
            errors.append(f"Invalid file type: {file.filename}")
            continue
            
        try:
        if file and allowed_file(file.filename):
            filename = file.filename
            file_ext = filename.rsplit('.', 1)[1].lower()
            file_id = str(uuid.uuid4())
            filename = f"{file_id}.pdf"
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            
            file.save(filepath)
            
            if not os.path.exists(filepath):
                errors.append(f"Failed to save: {file.filename}")
                continue
                
            # Process document
            text = extract_text_from_pdf(filepath)
            chunks, images = extract_headings_and_content(filepath=filepath)
            chunks = chunk_text(text) if text else []
            chunk_embeddings = embedding_model.encode(chunks).tolist() if chunks else []  # Precompute embeddings
            
            document_chunks[file_id] = {
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}.{file_ext}")
            print(f"Processing file: {filename}")
            file.save(file_path)
            file_size = os.path.getsize(file_path)
            text, headings, images = "", [], []
            if file_ext == 'pdf':
                text, headings = extract_text_from_pdf(file_path)
                images = extract_images_from_pdf(file_path)
            elif file_ext == 'pptx':
                text, headings = extract_text_from_pptx(file_path)
                images = extract_images_from_pptx(file_path)
            elif file_ext == 'docx':
                text, headings = extract_text_from_docx(file_path)
            elif file_ext == 'txt':
                text, headings = extract_text_from_txt(file_path)
            chunks = chunk_text(text)
            embeddings = generate_embeddings(chunks)
            document = {
                "id": file_id,
                "name": file.filename,
                "size": os.path.getsize(filepath),
                "name": filename,
                "size": file_size,
                "chunks": chunks,
                "chunk_embeddings": chunk_embeddings,
                "images": images,
                "upload_time": time.time()
            }
            
            uploaded_docs.append({
                "id": file_id,
                "name": file.filename,
                "size": os.path.getsize(filepath),
                "num_chunks": len(chunks),
                "embeddings": embeddings.tolist(),
                "headings": headings,
                "images": images
            })
            
        except Exception as e:
            logger.error(f"Error processing {file.filename}: {e}")
            errors.append(f"Failed to process {file.filename}")
            continue
    
    if uploaded_docs:
        with open(DATA_FILE, 'w') as f:
            json.dump(document_chunks, f)
            
        response = {
            "success": True,
            "message": f"Processed {len(uploaded_docs)} files successfully",
            "documents": uploaded_docs
        }
        
        if errors:
            response["warnings"] = errors
            
        return jsonify(response)
    
    return jsonify({
        "success": False,
        "message": "No valid files processed",
        "errors": errors
    }), 400
            }
            document_data["documents"].append(document)
            uploaded_documents.append({"id": file_id, "name": filename, "size": file_size, "images": images})
            print(f"File {filename} processed successfully")
    save_document_data(document_data)
    print(f"Uploaded {len(uploaded_documents)} documents")
    return jsonify({"success": True, "documents": uploaded_documents})

@app.route('/chat', methods=['POST'])
def handle_chat():
    log_debug("Chat request received")
    
    if not request.is_json:
        return jsonify({"success": False, "message": "JSON required"}), 400
    
def chat():
    print("Received chat query")
    data = request.get_json()
    user_message = data.get('message', '').strip()
    document_ids = data.get('documents', [])
    
    if not user_message:
        return jsonify({"success": False, "message": "Message required"}), 400
    
    if not data or 'message' not in data:
        print("No message provided in chat query")
        return jsonify({"success": False, "message": "No message provided"}), 400
    message = data['message']
    doc_ids = data.get('documents', [])
    relevant_chunks = []
    relevant_images = []
    print(f"Processing query: {message}")
    if doc_ids:
        print("Generating query embedding...")
        query_embedding = model.encode([message], convert_to_tensor=False)[0]
        print("Query embedding generated")
        for doc in document_data["documents"]:
            if doc["id"] in doc_ids:
                print(f"Searching document: {doc['name']}")
                embeddings = np.array(doc["embeddings"])
                similarities = cosine_similarity([query_embedding], embeddings)[0]
                top_indices = np.argsort(similarities)[-3:][::-1]
                print(f"Found {len(top_indices)} relevant chunks")
                for idx in top_indices:
                    if similarities[idx] > 0.2:
                        relevant_chunks.append(doc["chunks"][idx])
                        for image in doc["images"]:
                            if ("page" in image and image["page"] == idx + 1) or ("slide" in image and image["slide"] == idx + 1):
                                relevant_images.append({"url": f"/images/{image['filename']}", "heading": image["heading"], "similarity": similarities[idx]})
    context = "\n".join(relevant_chunks) if relevant_chunks else "No relevant document content found."
    prompt = f"Context: {context}\n\nQuestion: {message}\nAnswer:"
    print(f"Sending prompt to LLM with {len(relevant_chunks)} chunks and {len(relevant_images)} images")
    try:
        # Retrieve document chunks, embeddings, and images
        all_chunks = []
        all_chunk_embeddings = []
        all_images = []
        for doc_id in document_ids:
            if doc_id in document_chunks:
                all_chunks.extend(document_chunks[doc_id]['chunks'])
                all_chunk_embeddings.extend(document_chunks[doc_id].get('chunk_embeddings', []))
                all_images.extend(document_chunks[doc_id].get('images', []))
            else:
                log_debug(f"Document ID not found: {doc_id}")
        
        if not all_chunks:
            log_debug("No document chunks available - using general knowledge")
            response = query_llama(user_message)
            return jsonify({"success": True, "response": response})
        
        # Compute query embedding once
        search_embedding = embedding_model.encode([user_message])
        
        # Extract and analyze headings
        combined_text = "\n".join(chunk for chunk in all_chunks)
        log_debug(f"Combined text sample:\n{combined_text[:1000]}...")
        
        heading_content = extract_headings_and_content(raw_text=combined_text)[0]
        headings = list(heading_content.keys())
        log_debug(f"Extracted headings:\n{'- ' + '\n- '.join(headings)}")
        
        # Keyword extraction
        prompt = f"""Analyze this query and extract key terms that should appear in the answer.
Return ONLY a comma-separated list of 5-7 keywords/phrases.

Query: {user_message}"""
        
        keyword_string = query_llama(prompt)
        keywords = [k.strip() for k in keyword_string.split(",")] if keyword_string else []
        search_terms = f"{user_message} {' '.join(keywords)}"
        
        # Heading-based retrieval
        relevant_excerpts = []
        relevant_images = []
        relevant_headings = []
        if headings:
            log_debug("Attempting heading-based retrieval...")
            heading_embeddings = embedding_model.encode(headings)
            
            similarities = cosine_similarity(search_embedding, heading_embeddings)[0]
            heading_scores = sorted(zip(headings, similarities), 
                                  key=lambda x: x[1], reverse=True)
            
            relevant_headings = [h for h, s in heading_scores 
                               if s > HEADING_SIMILARITY_THRESHOLD][:3]
            
            if relevant_headings:
                log_debug(f"Found relevant headings: {relevant_headings}")
                for heading in relevant_headings:
                    content = heading_content.get(heading, "")
                    if content:
                        relevant_excerpts.append(f"## {heading}\n{content[:500]}")
                    # Find images associated with this heading
                    for img in all_images:
                        if not all(key in img for key in ['filename', 'heading']):
                            log_debug(f"Skipping invalid image: {img}")
                            continue
                        if img["heading"] == heading:
                            try:
                                similarity = similarities[headings.index(heading)] if heading in headings else 0
                                relevant_images.append({
                                    "url": f"/images/{img['filename']}",
                                    "heading": img["heading"],
                                    "similarity": float(similarity)  # Ensure JSON-serializable
                                })
                            except ValueError as e:
                                log_debug(f"Error finding heading index for {heading}: {e}")
                                relevant_images.append({
                                    "url": f"/images/{img['filename']}",
                                    "heading": img["heading"],
                                    "similarity": 0.0
                                })
        
        # Fallback to chunk-based retrieval if needed
        if not relevant_excerpts:
            log_debug("Using chunk-based retrieval fallback...")
            relevant_chunks = get_relevant_chunks(search_terms, all_chunks, all_chunk_embeddings)
            relevant_excerpts = [f"### Excerpt\n{chunk[0][:500]}" for chunk in relevant_chunks]
        
        # Generate response
        context = "\n\n".join(relevant_excerpts)[:5000]
        prompt = f"""Answer the question based on these document sections:

Question: {user_message}

Document Context:
{context}

Provide a comprehensive answer that:
1. Directly addresses the question
2. References relevant sections
3. Synthesizes information when needed
4. Is accurate and concise
5. Remember that your response is directly displayed on my website related to RAG so provide response accordingly"""

        response = query_llama(prompt)
        log_debug(f"Returning response with {len(relevant_images)} images")
        return jsonify({
            "success": True,
            "response": response,
            "used_headings": relevant_headings,
            "context": context,
            "images": relevant_images
        response = requests.post('http://localhost:11434/api/generate', json={
            "model": "llama3",
            "prompt": prompt,
            "stream": False
        })
        
        response_data = response.json()
        if response.status_code != 200 or 'response' not in response_data:
            print("Error from LLM API")
            raise Exception("Error from LLM API")
        print("Received LLM response")
        return jsonify({"success": True, "response": response_data["response"], "images": relevant_images})
    except Exception as e:
        logger.error(f"Chat error: {str(e)}", exc_info=True)
        return jsonify({
            "success": False,
            "message": "Processing error",
            "response": f"Sorry, I encountered an error: {str(e)}"
        }), 500
        print(f"Chat query failed: {str(e)}")
        return jsonify({"success": False, "message": str(e)}), 500

@app.route('/documents', methods=['GET'])
def list_documents():
    documents = []
    for doc_id, doc_data in document_chunks.items():
        documents.append({
            "id": doc_id,
            "name": doc_data.get("name", "Unknown"),
            "size": doc_data.get("size", 0),
            "num_chunks": len(doc_data.get("chunks", [])),
            "images": doc_data.get("images", [])
        })
    return jsonify({"success": True, "documents": documents})
def get_documents():
    print("Fetching document list")
    return jsonify({"documents": [{"id": doc["id"], "name": doc["name"], "size": doc["size"]} for doc in document_data["documents"]]})

@app.route('/images/<filename>')
def serve_image(filename):
    print(f"Serving image: {filename}")
    return send_from_directory(app.config['IMAGE_FOLDER'], filename)

@app.route('/', methods=['GET'])
def index():
    print("Serving index.html")
    return render_template('index.html')

if __name__ == '__main__':
    logger.info("Starting application")
    print("Starting Flask application...")
    app.run(host='0.0.0.0', port=5000, debug=True)
